"""
WPS MCP Server
Enables AI agents to create, edit, and manage WPS documents via natural language.
Note: WPS uses standard Office formats (.docx, .xlsx, .pptx),
so this server works with both WPS and Microsoft Office files.
"""

import os
from pathlib import Path
from typing import Optional

from mcp.server.fastmcp import FastMCP

# Initialize FastMCP server
mcp = FastMCP("wps")

# Default workspace
WORKSPACE = os.environ.get("WPS_WORKSPACE", str(Path.home() / "wps_projects"))


def _ensure_workspace() -> str:
    """Ensure workspace directory exists."""
    os.makedirs(WORKSPACE, exist_ok=True)
    return WORKSPACE


# ============================================================
# Document Tools (Word)
# ============================================================

@mcp.tool()
async def create_document(filename: str, title: str = "Untitled") -> str:
    """
    Create a new WPS Word document.

    Args:
        filename: Name of the document (without extension)
        title: Document title

    Returns:
        Path to created document
    """
    from docx import Document

    _ensure_workspace()
    doc_path = os.path.join(WORKSPACE, f"{filename}.docx")

    doc = Document()
    doc.add_heading(title, 0)
    doc.save(doc_path)

    return f"Created: {doc_path}"


@mcp.tool()
async def add_heading(filename: str, text: str, level: int = 1) -> str:
    """
    Add a heading to an existing document.

    Args:
        filename: Document filename
        text: Heading text
        level: Heading level (1-9, default 1)

    Returns:
        Confirmation message
    """
    from docx import Document

    doc_path = os.path.join(WORKSPACE, f"{filename}.docx")

    if not os.path.exists(doc_path):
        return f"Error: Document '{filename}' not found"

    doc = Document(doc_path)
    doc.add_heading(text, level)
    doc.save(doc_path)

    return f"Added heading '{text}' to {filename}"


@mcp.tool()
async def add_paragraph(filename: str, text: str, style: str = "Normal") -> str:
    """
    Add a paragraph to an existing document.

    Args:
        filename: Document filename
        text: Paragraph text
        style: Paragraph style (Normal, Quote, etc.)

    Returns:
        Confirmation message
    """
    from docx import Document

    doc_path = os.path.join(WORKSPACE, f"{filename}.docx")

    if not os.path.exists(doc_path):
        return f"Error: Document '{filename}' not found"

    doc = Document(doc_path)
    para = doc.add_paragraph(text)
    if style != "Normal":
        para.style = style

    doc.save(doc_path)
    return f"Added paragraph to {filename}"


@mcp.tool()
async def read_document(filename: str) -> str:
    """
    Read and return document content.

    Args:
        filename: Document filename

    Returns:
        Document content as text
    """
    from docx import Document

    doc_path = os.path.join(WORKSPACE, f"{filename}.docx")

    if not os.path.exists(doc_path):
        return f"Error: Document '{filename}' not found"

    doc = Document(doc_path)
    content = []
    for para in doc.paragraphs:
        if para.text.strip():
            content.append(para.text)

    return "\n".join(content) if content else "Document is empty"


# ============================================================
# Spreadsheet Tools (Excel)
# ============================================================

@mcp.tool()
async def create_spreadsheet(filename: str, sheet_name: str = "Sheet1") -> str:
    """
    Create a new WPS Excel spreadsheet.

    Args:
        filename: Name of the spreadsheet (without extension)
        sheet_name: Name of the first sheet

    Returns:
        Path to created spreadsheet
    """
    from openpyxl import Workbook

    _ensure_workspace()
    xlsx_path = os.path.join(WORKSPACE, f"{filename}.xlsx")

    wb = Workbook()
    wb.active.title = sheet_name
    wb.save(xlsx_path)

    return f"Created: {xlsx_path}"


@mcp.tool()
async def write_cell(
    filename: str, cell: str, value: str, sheet: str = "Sheet1"
) -> str:
    """
    Write a value to a cell.

    Args:
        filename: Spreadsheet filename
        cell: Cell address (e.g., "A1", "B2")
        value: Value to write
        sheet: Sheet name

    Returns:
        Confirmation message
    """
    from openpyxl import load_workbook

    xlsx_path = os.path.join(WORKSPACE, f"{filename}.xlsx")

    if not os.path.exists(xlsx_path):
        return f"Error: Spreadsheet '{filename}' not found"

    wb = load_workbook(xlsx_path)

    if sheet not in wb.sheetnames:
        wb.create_sheet(sheet)

    ws = wb[sheet]
    ws[cell] = value
    wb.save(xlsx_path)

    return f"Wrote '{value}' to {cell} in {filename}"


@mcp.tool()
async def read_cell(filename: str, cell: str, sheet: str = "Sheet1") -> str:
    """
    Read a cell value.

    Args:
        filename: Spreadsheet filename
        cell: Cell address (e.g., "A1")
        sheet: Sheet name

    Returns:
        Cell value
    """
    from openpyxl import load_workbook

    xlsx_path = os.path.join(WORKSPACE, f"{filename}.xlsx")

    if not os.path.exists(xlsx_path):
        return f"Error: Spreadsheet '{filename}' not found"

    wb = load_workbook(xlsx_path, data_only=True)

    if sheet not in wb.sheetnames:
        return f"Error: Sheet '{sheet}' not found"

    ws = wb[sheet]
    value = ws[cell].value

    return f"{cell}: {value}" if value is not None else f"{cell}: (empty)"


@mcp.tool()
async def add_formula(filename: str, cell: str, formula: str, sheet: str = "Sheet1") -> str:
    """
    Add a formula to a cell.

    Args:
        filename: Spreadsheet filename
        cell: Cell address (e.g., "B1")
        formula: Formula (e.g., "=SUM(A1:A10)")
        sheet: Sheet name

    Returns:
        Confirmation message
    """
    from openpyxl import load_workbook

    xlsx_path = os.path.join(WORKSPACE, f"{filename}.xlsx")

    if not os.path.exists(xlsx_path):
        return f"Error: Spreadsheet '{filename}' not found"

    wb = load_workbook(xlsx_path)

    if sheet not in wb.sheetnames:
        wb.create_sheet(sheet)

    ws = wb[sheet]
    ws[cell] = formula
    wb.save(xlsx_path)

    return f"Added formula to {cell} in {filename}"


# ============================================================
# Presentation Tools (PowerPoint)
# ============================================================

@mcp.tool()
async def create_presentation(filename: str, title: str = "Presentation") -> str:
    """
    Create a new WPS PowerPoint presentation.

    Args:
        filename: Name of the presentation (without extension)
        title: Presentation title

    Returns:
        Path to created presentation
    """
    from pptx import Presentation

    _ensure_workspace()
    pptx_path = os.path.join(WORKSPACE, f"{filename}.pptx")

    prs = Presentation()
    prs.slide_width = 9144000  # 10 inches
    prs.slide_height = 5143500  # 5.625 inches

    # Add title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    if title_slide.shapes.title:
        title_slide.shapes.title.text = title

    prs.save(pptx_path)

    return f"Created: {pptx_path}"


@mcp.tool()
async def add_slide(filename: str, layout_index: int = 1) -> str:
    """
    Add a new slide to presentation.

    Args:
        filename: Presentation filename
        layout_index: Layout index (0=title, 1=blank, etc.)

    Returns:
        Confirmation message
    """
    from pptx import Presentation

    pptx_path = os.path.join(WORKSPACE, f"{filename}.pptx")

    if not os.path.exists(pptx_path):
        return f"Error: Presentation '{filename}' not found"

    prs = Presentation(pptx_path)

    if layout_index >= len(prs.slide_layouts):
        layout_index = 1

    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    prs.save(pptx_path)

    slide_num = len(prs.slides)
    return f"Added slide {slide_num} to {filename}"


@mcp.tool()
async def add_text_to_slide(
    filename: str, text: str, left: float = 1, top: float = 1,
    width: float = 8, height: float = 1
) -> str:
    """
    Add text box to a slide.

    Args:
        filename: Presentation filename
        text: Text content
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches

    Returns:
        Confirmation message
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    pptx_path = os.path.join(WORKSPACE, f"{filename}.pptx")

    if not os.path.exists(pptx_path):
        return f"Error: Presentation '{filename}' not found"

    prs = Presentation(pptx_path)

    if not prs.slides:
        return f"Error: No slides in {filename}"

    slide = prs.slides[-1]
    left_in = Inches(left)
    top_in = Inches(top)
    width_in = Inches(width)
    height_in = Inches(height)

    textbox = slide.shapes.add_textbox(left_in, top_in, width_in, height_in)
    textbox.text_frame.text = text

    prs.save(pptx_path)
    return f"Added text to {filename}"


@mcp.tool()
async def set_slide_layout(filename: str, layout_index: int = 1) -> str:
    """
    Change slide layout.

    Args:
        filename: Presentation filename
        layout_index: Layout index (0=title, 1=blank, etc.)

    Returns:
        Confirmation message
    """
    from pptx import Presentation

    pptx_path = os.path.join(WORKSPACE, f"{filename}.pptx")

    if not os.path.exists(pptx_path):
        return f"Error: Presentation '{filename}' not found"

    prs = Presentation(pptx_path)

    if not prs.slides:
        return f"Error: No slides in {filename}"

    if layout_index >= len(prs.slide_layouts):
        return f"Error: Layout index {layout_index} out of range"

    # Get layout from presentation
    desired_layout = prs.slide_layouts[layout_index]
    slide = prs.slides[-1]

    # Copy placeholders from layout to slide
    for placeholder in slide.placeholders:
        placeholder.element.getparent().remove(placeholder.element)

    for shape in desired_layout.shapes:
        if shape.has_text_frame:
            el = shape.element
            new_el = type(el)(el)
            slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    prs.save(pptx_path)
    return f"Set layout {layout_index} for latest slide in {filename}"


# ============================================================
# Utility
# ============================================================

@mcp.tool()
async def list_files() -> str:
    """
    List all WPS files in workspace.

    Returns:
        List of files
    """
    _ensure_workspace()
    files = []
    for f in os.listdir(WORKSPACE):
        if f.endswith(('.docx', '.xlsx', '.pptx')):
            files.append(f)
    return "\n".join(files) if files else "No files found"


@mcp.tool()
async def get_workspace() -> str:
    """
    Get the current workspace path.

    Returns:
        Workspace path
    """
    return _ensure_workspace()


if __name__ == "__main__":
    mcp.run()
